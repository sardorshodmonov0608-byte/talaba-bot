import os
import sqlite3
import logging
import random
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import CommandStart
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup

# PPTX dizayni uchun kerakli modullar
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Word hujjati uchun
from docx import Document

logging.basicConfig(level=logging.INFO)
DB_NAME = "talaba_premium.db"

# --- INTELLEKTUAL AI MATN GENERATORI (MOCK AI) ---
def generate_ai_content(topic, mode="essay"):
    clean_topic = topic.strip().capitalize()
    if mode == "essay":
        return {
            "title": clean_topic,
            "intro": f"Bugungi globallashuv va shiddatli rivojlanish davrida '{clean_topic}' mavzusi nafaqat ilmiy doiralarda, balki jamiyatning har bir qatlamida dolzarb ahamiyat kasb etmoqda. Ushbu masalaning tub mohiyatini anglash kelajak poydevorini mustahkamlashga xizmat qiladi.",
            "body_1": f"Tahlillarga ko'ra, '{clean_topic}' yo'nalishidagi asosiy muammolar va yechimlar tizimli yondashuvni talab etadi. Bu jarayonda zamonaviy texnologiyalar, innovatsion metodlar hamda xorijiy tajribalardan unumli foydalanish yuqori samara berishi amaliyotda isbotlangan.",
            "body_2": "Muammoning ijtimoiy-iqtisodiy va huquqiy jihatlariga to'xtaladigan bo'lsak, har bir bosqichni individual tahlil qilish lozim. Yosh avlod va soha mutaxassislarining bu boradagi ko'nikmalarini oshirish eng asosiy drayver hisoblanadi.",
            "conclusion": f"Xulosa qilib aytganda, '{clean_topic}' sohasini rivojlantirish kompleks chora-tadbirlarga bog'liq. Bugun amalga oshirilayotgan islohotlar ertangi kunning yorqin natijalariga zamin yaratadi deb ishonch bilan ayta olamiz."
        }
    else:
        return [
            {"title": clean_topic, "content": "Talaba Bot AI tizimi tomonidan tayyorlangan taqdimot\n\nToshkent — 2026"},
            {"title": "1. Kirish va Dolzarbligi", "content": f"• '{clean_topic}' tushunchasining kelib chiqishi\n• Bugungi kundagi ijtimoiy va ilmiy ahamiyati\n• Taqdimotning asosiy maqsad va vazifalari"},
            {"title": "2. Asosiy muammolar", "content": "• Sohadagi mavjud kamchiliklar va to'siqlar\n• Statistik ma'lumotlar va real tahlillar\n• Salbiy omillarning ta'sir doirasi"},
            {"title": "3. Innovatsion Yechimlar", "content": "• Muammolarni bartaraf etishning zamonaviy usullari\n• Raqamlashtirish va sun'iy intellekt integratsiyasi\n• Kutilayotgan samaradorlik ko'rsatkichlari"},
            {"title": "4. Xulosa va Tavsiyalar", "content": f"• '{clean_topic}' bo'yicha yakuniy xulosalar\n• Amaliyotga joriy etish tavsiyalari\n• E'tiboringiz uchun rahmat!"}
        ]

# --- BAZA TIZIMI ---
def init_db():
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS users (
            user_id INTEGER PRIMARY KEY,
            free_attempts INTEGER DEFAULT 3,
            is_vip INTEGER DEFAULT 0
        )
    ''')
    conn.commit()
    conn.close()

def get_or_create_user(user_id):
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    cursor.execute("SELECT free_attempts, is_vip FROM users WHERE user_id = ?", (user_id,))
    user = cursor.fetchone()
    if not user:
        cursor.execute("INSERT INTO users (user_id, free_attempts, is_vip) VALUES (?, 3, 0)", (user_id,))
        conn.commit()
        user = (3, 0)
    conn.close()
    return {"free_attempts": user[0], "is_vip": user[1]}

def decrease_attempt(user_id):
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    cursor.execute("UPDATE users SET free_attempts = free_attempts - 1 WHERE user_id = ?", (user_id,))
    conn.commit()
    conn.close()

def set_vip_status(user_id):
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    cursor.execute("UPDATE users SET is_vip = 1, free_attempts = 0 WHERE user_id = ?", (user_id,))
    conn.commit()
    conn.close()

init_db()class StudentBotStates(StatesGroup):
    waiting_for_slide_topic = State()
    waiting_for_essay_topic = State()

def main_keyboard():
    kb = [
        [types.KeyboardButton(text="📊 Premium Slayd Yaratish"), types.KeyboardButton(text="📝 Ilmiy Esse Yozish")],
        [types.KeyboardButton(text="💎 VIP Obuna (Cheksiz)"), types.KeyboardButton(text="👤 Profilim")]
    ]
    return types.ReplyKeyboardMarkup(keyboard=kb, resize_keyboard=True)

dp = Dispatcher()

@dp.message(CommandStart())
async def start_cmd(message: types.Message):
    user_data = get_or_create_user(message.from_user.id)
    welcome = (
        f"🚀 Salom {message.from_user.full_name}! O'zbekistondagi eng kuchli Talaba Botga xush kelibsiz!\n\n"
        f"Men sun'iy intellekt yordamida bir necha soniyada dizaynga ega Slaydlar va original Esselar yaratib bera olaman.\n\n"
        f"🎁 Sizda {user_data['free_attempts']} ta bepul premium imkoniyat bor!"
    )
    await message.answer(welcome, reply_markup=main_keyboard(), parse_mode="Markdown")

@dp.message(F.text == "👤 Profilim")
async def show_profile(message: types.Message):
    user_data = get_or_create_user(message.from_user.id)
    status = "💎 VIP Status (Cheksiz)" if user_data['is_vip'] else f"📉 Bepul urinishlaringiz: {user_data['free_attempts']} ta"
    await message.answer(
        f"👤 Foydalanuvchi profili:\n\n"
        f"🆔 Sizning ID: {message.from_user.id}\n"
        f"⚡ Status: {status}",
        parse_mode="Markdown"
    )

async def check_user_access(message: types.Message) -> bool:
    user_data = get_or_create_user(message.from_user.id)
    if user_data['is_vip'] == 1 or user_data['free_attempts'] > 0:
        return True
    await message.answer(
        "❌ Kechirasiz, bepul urinishlaringiz tugadi!\n\n"
        "Botdan cheksiz foydalanishni davom ettirish uchun VIP obunani faollashtiring.",
        reply_markup=main_keyboard()
    )
    return False

@dp.message(F.text == "📊 Premium Slayd Yaratish")
async def ask_slide_topic(message: types.Message, state: FSMContext):
    if not await check_user_access(message): return
    await message.answer("✨ Slayd mavzusini kiriting:\n*(Masalan: Sun'iy intellekt va kelajak)*")
    await state.set_state(StudentBotStates.waiting_for_slide_topic)

@dp.message(StudentBotStates.waiting_for_slide_topic)
async def create_premium_slide(message: types.Message, state: FSMContext):
    topic = message.text
    user_id = message.from_user.id
    user_data = get_or_create_user(user_id)
    await message.answer("🔄 AI mavzuni tahlil qilmoqda va dizayn chizmoqda...")
    slides_data = generate_ai_content(topic, mode="slide")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    BG_COLOR = RGBColor(11, 19, 43)       
    TEXT_COLOR = RGBColor(255, 255, 255)  
    ACCENT_COLOR = RGBColor(0, 180, 216)  

    for i, s_data in enumerate(slides_data):
        blank_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_layout)
        bg_shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_COLOR
        bg_shape.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.333), Inches(1.5))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = s_data["title"]
        p_title.font.size = Pt(40 if i == 0 else 34)
        p_title.font.bold = True
        p_title.font.color.rgb = ACCENT_COLOR
        if i == 0:
            p_title.alignment = PP_ALIGN.CENTER
            title_box.top = Inches(2.2)

        content_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(4))tf_content = content_box.text_frame
        tf_content.word_wrap = True
        p_content = tf_content.paragraphs[0]
        p_content.text = s_data["content"]
        p_content.font.size = Pt(22)
        p_content.font.color.rgb = TEXT_COLOR
        if i == 0:
            p_content.alignment = PP_ALIGN.CENTER
            content_box.top = Inches(4.2)

    file_name = f"{user_id}_slayd.pptx"
    prs.save(file_name)
    if not user_data['is_vip']: decrease_attempt(user_id)
    await message.answer_document(types.FSInputFile(file_name), caption=f"🎬 🔥 '{topic}' mavzusidagi slayd tayyor!")
    os.remove(file_name)
    await state.clear()

@dp.message(F.text == "📝 Ilmiy Esse Yozish")
async def ask_essay_topic(message: types.Message, state: FSMContext):
    if not await check_user_access(message): return
    await message.answer("✍️ Esse mavzusini kiriting:\n*(Masalan: Kitobxonlik va inson ma'naviyati)*")
    await state.set_state(StudentBotStates.waiting_for_essay_topic)

@dp.message(StudentBotStates.waiting_for_essay_topic)
async def create_premium_essay(message: types.Message, state: FSMContext):
    topic = message.text
    user_id = message.from_user.id
    user_data = get_or_create_user(user_id)
    await message.answer("✍️ AI esse yozmoqda...")
    essay = generate_ai_content(topic, mode="essay")

    doc = Document()
    doc.add_heading(essay["title"], level=0)
    doc.add_heading("1. Kirish qismi", level=1)
    doc.add_paragraph(essay["intro"])
    doc.add_heading("2. Tahliliy va asosiy qism", level=1)
    doc.add_paragraph(essay["body_1"])
    doc.add_paragraph(essay["body_2"])
    doc.add_heading("3. Yakuniy xulosa", level=1)
    doc.add_paragraph(essay["conclusion"])

    file_name = f"{user_id}_esse.docx"
    doc.save(file_name)
    if not user_data['is_vip']: decrease_attempt(user_id)
    await message.answer_document(types.FSInputFile(file_name), caption=f"📄 '{topic}' mavzusidagi esse tayyor!")
    os.remove(file_name)
    await state.clear()

@dp.message(F.text == "💎 VIP Obuna (Cheksiz)")
async def send_payment_info(message: types.Message):
    user_data = get_or_create_user(message.from_user.id)
    if user_data['is_vip']:
        await message.answer("🎉 Sizda allaqachon VIP obuna faol!")
        return
    text = (
        "💎 VIP Obuna — Cheksiz Imkoniyatlar!\n\n"
        "💰 1 oylik obuna narxi: 25 000 so'm\n\n"
        "To'lovni Mavrid orqali o'tkazing:\n"
        "💳 Karta raqam: 9860012126025795\n"
        "👤 Qabul qiluvchi: shodmonov sardor\n\n"
        "⚠️ To'lovdan so'ng chekni (rasmni) shu yerga jo'nating!"
    )
    await message.answer(text, parse_mode="Markdown")

@dp.message(F.photo)
async def handle_screenshot(message: types.Message):
    user_id = message.from_user.id
    username = f"@{message.from_user.username}" if message.from_user.username else "Yo'q"
    kb = [[types.InlineKeyboardButton(text="✅ Tasdiqlash", callback_data=f"accept_{user_id}"),
           types.InlineKeyboardButton(text="❌ Rad etish", callback_data=f"reject_{user_id}")]]
    await bot.send_photo(chat_id=ADMIN_ID, photo=message.photo[-1].file_id, 
                         caption=f"🔔 To'lov cheki!\n👤 Kimdan: {message.from_user.full_name}\n🆔 ID: {user_id}", 
                         reply_markup=types.InlineKeyboardMarkup(inline_keyboard=kb))
    await message.answer("🔄 Chek adminga yuborildi. Tez orada VIP faollashadi!")

@dp.callback_query(F.data.startswith("accept_"))
async def approve_vip(callback: types.CallbackQuery):
    target_user_id = int(callback.data.split("_")[1])
    set_vip_status(target_user_id)
    try:
        await bot.send_message(chat_id=target_user_id, text="🎉 VIP obuna berildi!")
    except: pass
    await callback.message.edit_caption(caption=callback.message.caption + "\n\n🟢 VIP FAOLLASHTIRILDI")@dp.callback_query(F.data.startswith("reject_"))
async def reject_vip(callback: types.CallbackQuery):
    target_user_id = int(callback.data.split("_")[1])
    try:
        await bot.send_message(chat_id=target_user_id, text="❌ To'lov chekingiz tasdiqlanmadi.")
    except: pass
    await callback.message.edit_caption(caption=callback.message.caption + "\n\n🔴 RAD ETILDI")

# ====================================================================
# 👇 ENGI PASTI — MANA SHU YERGA O'ZINGIZNIKINI YOZING 👇
# ====================================================================

BOT_TOKEN = "8802435742:AAGW4o9UjsRAtXvsHMj2mlrApBWg8pPTuts"   # 👈 O'rniga BotFather bergan tokenni qo'ying
ADMIN_ID = 6141302755                    # 👈 O'rniga userinfobot bergan ID raqamingizni yozing

# ====================================================================

if name == "main":
    import asyncio
    bot = Bot(token=BOT_TOKEN) # Botni yangi token bilan yuklash
    asyncio.run(dp.start_polling(bot))
